#!/usr/bin/env python3
"""Build the AARP x Parigrado pitch deck (<=10 slides) with python-pptx."""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# Palette
AARP_RED = RGBColor(0xC8, 0x10, 0x2E)
CHARCOAL = RGBColor(0x2B, 0x2B, 0x2B)
GRAY = RGBColor(0x60, 0x66, 0x6E)
LIGHT = RGBColor(0xF4, 0xF5, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color=WHITE):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def setpara(p, text, size, color, bold=False, align=PP_ALIGN.LEFT,
            font=FONT, space_after=6, space_before=0, italic=False):
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return r


def eyebrow(s, text, l=Inches(0.9), t=Inches(0.55), w=Inches(11.5)):
    _, tf = box(s, l, t, w, Inches(0.4))
    setpara(tf.paragraphs[0], text.upper(), 12.5, AARP_RED, bold=True)
    # letter spacing look via caps only
    return tf


def headline(s, text, t=Inches(0.95), size=38, l=Inches(0.9), w=Inches(11.5), color=CHARCOAL):
    _, tf = box(s, l, t, w, Inches(1.3))
    setpara(tf.paragraphs[0], text, size, color, bold=True)
    return tf


def redbar(s, l=Inches(0.92), t=Inches(1.72), w=Inches(0.85), h=Inches(0.06)):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = AARP_RED
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def rect(s, l, t, w, h, fill=LIGHT, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sp.adjustments[0] = 0.06
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def bullets(tf, items, size=17, color=CHARCOAL, gap=12, bold_lead=True):
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.alignment = PP_ALIGN.LEFT
        # red dot
        dot = p.add_run()
        dot.text = "•  "
        dot.font.size = Pt(size)
        dot.font.bold = True
        dot.font.color.rgb = AARP_RED
        dot.font.name = FONT
        if isinstance(it, tuple):
            lead, rest = it
            r1 = p.add_run(); r1.text = lead + " "
            r1.font.size = Pt(size); r1.font.bold = True
            r1.font.color.rgb = CHARCOAL; r1.font.name = FONT
            r2 = p.add_run(); r2.text = rest
            r2.font.size = Pt(size); r2.font.color.rgb = GRAY; r2.font.name = FONT
        else:
            r = p.add_run(); r.text = it
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT


def disclaimer(s, text="Partnership concept — not an official AARP product."):
    _, tf = box(s, Inches(0.9), Inches(7.02), Inches(11.5), Inches(0.4))
    setpara(tf.paragraphs[0], text, 9.5, GRAY, italic=True)


# ---------------- Slide 1: Title ----------------
s = slide(); bg(s, WHITE)
# left red accent band
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SH)
band.fill.solid(); band.fill.fore_color.rgb = AARP_RED; band.line.fill.background(); band.shadow.inherit = False
_, tf = box(s, Inches(1.0), Inches(0.7), Inches(11.5), Inches(0.4))
setpara(tf.paragraphs[0], "PARTNERSHIP CONCEPT  ·  CONFIDENTIAL", 12.5, AARP_RED, bold=True)
_, tf = box(s, Inches(1.0), Inches(2.35), Inches(11.5), Inches(2.0))
setpara(tf.paragraphs[0], "Parigrado × AARP", 60, CHARCOAL, bold=True)
p = tf.add_paragraph()
setpara(p, "Unbiased hospital quality for the 50+ community", 26, GRAY, space_before=6)
redbar(s, l=Inches(1.02), t=Inches(4.55), w=Inches(1.1), h=Inches(0.07))
_, tf = box(s, Inches(1.0), Inches(4.8), Inches(10.8), Inches(1.2))
setpara(tf.paragraphs[0],
        "Turning the government's own CMS and CDC data into clear, side-by-side hospital "
        "comparisons — no advertising, no sponsorships, no pay-to-rank lists.", 17, GRAY)
_, tf = box(s, Inches(1.0), Inches(6.9), Inches(11.5), Inches(0.4))
setpara(tf.paragraphs[0], "Not an official AARP product  ·  parigrado.com", 11, GRAY, italic=True)

# ---------------- Slide 2: Problem ----------------
s = slide(); bg(s)
eyebrow(s, "The Problem")
headline(s, "Choosing a hospital is high-stakes — and confusing")
redbar(s)
_, tf = box(s, Inches(0.9), Inches(2.2), Inches(6.4), Inches(4.2))
bullets(tf, [
    ("Consequential.", "Where you get care is one of the most important health decisions a 50+ patient makes."),
    ("Opaque.", "Quality varies widely between hospitals, but the differences are hard to see."),
    ("Ad-driven rankings.", "Many \u201cbest hospital\u201d lists are shaped by marketing, sponsorships, and pay-to-rank influence."),
    ("Hard to compare.", "Federal data exists, but it isn't presented in plain language for the people who need it."),
], size=18, gap=16)
# right stat card
rect(s, Inches(7.7), Inches(2.2), Inches(4.7), Inches(4.0), fill=LIGHT)
_, tf = box(s, Inches(8.05), Inches(2.55), Inches(4.0), Inches(3.4))
setpara(tf.paragraphs[0], "100M+", 46, AARP_RED, bold=True)
setpara(tf.add_paragraph(), "U.S. adults age 50+ facing these decisions", 15, GRAY, space_after=18)
setpara(tf.add_paragraph(), "$0", 46, CHARCOAL, bold=True)
setpara(tf.add_paragraph(), "Members should never pay to see unbiased quality data", 15, GRAY)
disclaimer(s)

# ---------------- Slide 3: Solution ----------------
s = slide(); bg(s)
eyebrow(s, "The Solution")
headline(s, "Parigrado: facts, not marketing")
redbar(s)
_, tf = box(s, Inches(0.9), Inches(2.1), Inches(11.4), Inches(0.9))
setpara(tf.paragraphs[0],
        "Search any Medicare-certified hospital and instantly see how it compares to its county, "
        "state, and national peers on the measures that matter to patients.", 18, GRAY)
_, tf = box(s, Inches(0.9), Inches(3.1), Inches(11.4), Inches(3.2))
bullets(tf, [
    ("Built on public federal data.", "Every score comes from CMS and CDC — presented as-is."),
    ("No pay-to-rank, ever.", "Nothing is for sale; no hospital can buy a higher ranking."),
    ("Made for a 50+ audience.", "Large-type, plain-language comparisons — not clinical jargon."),
    ("White-label ready.", "Re-skinned in AARP chrome and embedded via a simple link or iframe."),
], size=18, gap=16)
disclaimer(s)

# ---------------- Slide 4: Product snapshot (feature grid) ----------------
s = slide(); bg(s)
eyebrow(s, "Product Snapshot")
headline(s, "Everything members need, in one place")
redbar(s)
feats = [
    ("Peer benchmarking", "Compare any hospital against similar peers — county, state, national."),
    ("Historical trends", "Track how a hospital's performance moves over time."),
    ("PDF / CSV export", "Download a clean report in one click to keep or print."),
    ("Shareable links", "Send a live comparison to family or a physician."),
]
gx = [Inches(0.9), Inches(6.95)]
gy = [Inches(2.35), Inches(4.55)]
cw, ch = Inches(5.45), Inches(1.95)
for i, (title, desc) in enumerate(feats):
    l = gx[i % 2]; t = gy[i // 2]
    rect(s, l, t, cw, ch, fill=LIGHT)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Inches(0.12), ch)
    accent.fill.solid(); accent.fill.fore_color.rgb = AARP_RED
    accent.line.fill.background(); accent.shadow.inherit = False
    _, tf = box(s, l + Inches(0.4), t + Inches(0.28), cw - Inches(0.7), ch - Inches(0.5))
    setpara(tf.paragraphs[0], title, 20, CHARCOAL, bold=True)
    setpara(tf.add_paragraph(), desc, 15, GRAY, space_before=4)
disclaimer(s)

# ---------------- Slide 5: Market & Target Customer ----------------
s = slide(); bg(s)
eyebrow(s, "Market & Target Customer")
headline(s, "A large, aging, decision-heavy audience")
redbar(s)
stats = [("100M+", "U.S. adults age 50+"), ("~38M", "AARP members"),
         ("Growing", "Aging population + AgeTech demand")]
x = Inches(0.9)
for val, lab in stats:
    rect(s, x, Inches(2.15), Inches(3.7), Inches(1.7), fill=LIGHT)
    _, tf = box(s, x + Inches(0.25), Inches(2.35), Inches(3.2), Inches(1.3))
    setpara(tf.paragraphs[0], val, 34, AARP_RED, bold=True)
    setpara(tf.add_paragraph(), lab, 14, GRAY, space_before=2)
    x = Emu(x + Inches(3.9))
_, tf = box(s, Inches(0.9), Inches(4.25), Inches(11.4), Inches(2.4))
bullets(tf, [
    ("Primary customer.", "AARP and other trusted white-label / membership partners who serve the 50+ community."),
    ("End users.", "Members age 50+ (and their families) choosing where to receive hospital care."),
    ("Category tailwind.", "Rising interest in AgeTech and consumer healthcare-decision tools — framed directionally, not as invented TAM."),
], size=17, gap=13)
disclaimer(s)

# ---------------- Slide 6: Member value / Why AARP ----------------
s = slide(); bg(s)
eyebrow(s, "Member Value  ·  Why AARP")
headline(s, "On-mission, on-brand, and sticky")
redbar(s)
cards = [
    ("Mission fit", "Gives members facts, not marketing, on one of their most consequential decisions — squarely on AARP's consumer-empowerment mission."),
    ("Trust, protected", "Every score comes from CMS and CDC. Nothing is for sale, protecting AARP's brand equity around objectivity."),
    ("Engagement + retention", "A genuinely useful utility gives members a reason to return — measurable through sessions, comparisons, and repeat use."),
]
x = Inches(0.9); cw = Inches(3.75)
for title, desc in cards:
    rect(s, x, Inches(2.3), cw, Inches(3.9), fill=LIGHT)
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.3), cw, Inches(0.12))
    top.fill.solid(); top.fill.fore_color.rgb = AARP_RED; top.line.fill.background(); top.shadow.inherit = False
    _, tf = box(s, x + Inches(0.3), Inches(2.7), cw - Inches(0.6), Inches(3.3))
    setpara(tf.paragraphs[0], title, 20, CHARCOAL, bold=True, space_after=10)
    setpara(tf.add_paragraph(), desc, 15, GRAY)
    x = Emu(x + Inches(3.95))
disclaimer(s)

# ---------------- Slide 7: Partnership model + Revenue ----------------
s = slide(); bg(s)
eyebrow(s, "Partnership Model & Revenue")
headline(s, "White-label embed — low lift for AARP")
redbar(s)
# table-like rows
rows = [
    ("Embed / white-label", "AARP-branded Parigrado embedded via link or iframe.", "Preferred"),
    ("Licensed co-brand", "Deeper integration: AARP nav, member SSO, co-branded reports.", "Phase 2"),
    ("Pilot → national", "Prove value in 2–3 states, then scale nationally with SLAs.", "Recommended path"),
]
ty = Inches(2.25)
# header row
hdr = rect(s, Inches(0.9), ty, Inches(11.5), Inches(0.5), fill=CHARCOAL)
for txt, l, w in [("Model", Inches(1.05), Inches(3.0)),
                  ("What it is", Inches(4.1), Inches(5.6)),
                  ("Fit", Inches(9.9), Inches(2.3))]:
    _, tf = box(s, l, ty + Inches(0.06), w, Inches(0.4))
    setpara(tf.paragraphs[0], txt, 14, WHITE, bold=True)
ty = Emu(ty + Inches(0.5))
for i, (m, w, fit) in enumerate(rows):
    fill = WHITE if i % 2 == 0 else LIGHT
    rect(s, Inches(0.9), ty, Inches(11.5), Inches(0.85), fill=fill, line=RGBColor(0xE2,0xE4,0xE8))
    _, tf = box(s, Inches(1.05), ty + Inches(0.16), Inches(3.0), Inches(0.6))
    setpara(tf.paragraphs[0], m, 15, CHARCOAL, bold=True)
    _, tf = box(s, Inches(4.1), ty + Inches(0.16), Inches(5.6), Inches(0.6))
    setpara(tf.paragraphs[0], w, 13.5, GRAY)
    _, tf = box(s, Inches(9.9), ty + Inches(0.16), Inches(2.3), Inches(0.6))
    setpara(tf.paragraphs[0], fit, 14, AARP_RED, bold=True)
    ty = Emu(ty + Inches(0.85))
_, tf = box(s, Inches(0.9), Inches(6.25), Inches(11.5), Inches(0.7))
setpara(tf.paragraphs[0],
        "Revenue model: B2B — the partner licenses the tool; members always use it free. "
        "Detailed commercial terms are a pilot-first conversation.", 14.5, GRAY, italic=True)
disclaimer(s)

# ---------------- Slide 8: 90-day pilot + Traction ----------------
s = slide(); bg(s)
eyebrow(s, "The Ask  ·  90-Day Pilot & Traction")
headline(s, "A low-lift, fully-operated pilot")
redbar(s)
phases = [
    ("Weeks 1–2 · Stand up", "Configure the AARP white-label embed and disclaimers; align brand assets."),
    ("Weeks 3–8 · Soft launch", "Introduce to members in 2–3 high-50+ states; gather feedback."),
    ("Weeks 9–12 · Measure", "Review sessions, comparisons, return rate and sentiment; go / no-go readout."),
]
x = Inches(0.9); cw = Inches(3.75)
for title, desc in phases:
    rect(s, x, Inches(2.2), cw, Inches(2.2), fill=LIGHT)
    _, tf = box(s, x + Inches(0.28), Inches(2.4), cw - Inches(0.55), Inches(1.9))
    setpara(tf.paragraphs[0], title, 16.5, AARP_RED, bold=True, space_after=8)
    setpara(tf.add_paragraph(), desc, 14, GRAY)
    x = Emu(x + Inches(3.95))
_, tf = box(s, Inches(0.9), Inches(4.75), Inches(11.5), Inches(2.1))
setpara(tf.paragraphs[0], "Where we are today (honest, early-stage):", 15, CHARCOAL, bold=True, space_after=8)
bullets(tf, [
    ("Product is live", "at parigrado.com, running on a working CMS/CDC data pipeline."),
    ("AARP white-label demo", "is built and kept private — available on request for partner review."),
    ("Partnership-ready", "seeking a first pilot partner; no paying customers or investor raise to report yet."),
], size=14.5, gap=9)
disclaimer(s)

# ---------------- Slide 9: Team & Data trust ----------------
s = slide(); bg(s)
eyebrow(s, "Team & Data Trust")
headline(s, "Founder-led, built on public federal data")
redbar(s)
# left: team
rect(s, Inches(0.9), Inches(2.2), Inches(5.5), Inches(4.1), fill=LIGHT)
_, tf = box(s, Inches(1.2), Inches(2.5), Inches(5.0), Inches(3.6))
setpara(tf.paragraphs[0], "Team", 20, CHARCOAL, bold=True, space_after=10)
bullets(tf, [
    ("Founder-led.", "Led by Larry Kirschner, with a background in consumer-facing HOA and lending services (HOA Loan Services)."),
    ("Operator mindset.", "Experience building practical tools for everyday consumers, applied to healthcare decisions."),
    ("Lean & focused.", "Parigrado operates the product, data pipeline, and updates end to end."),
], size=14.5, gap=11)
# right: data trust
rect(s, Inches(6.7), Inches(2.2), Inches(5.7), Inches(4.1), fill=LIGHT)
_, tf = box(s, Inches(7.0), Inches(2.5), Inches(5.2), Inches(3.6))
setpara(tf.paragraphs[0], "Data trust & sourcing", 20, CHARCOAL, bold=True, space_after=10)
bullets(tf, [
    "CMS Hospital Compare",
    "CMS / CDC infection reporting",
    "CMS Hospital Readmissions Reduction Program",
], size=14.5, gap=8)
_, tf2 = box(s, Inches(7.0), Inches(5.15), Inches(5.2), Inches(1.1))
setpara(tf2.paragraphs[0],
        "Public data presented as-is, for informational purposes only — not medical advice "
        "and not a sole basis for any care decision.", 13, GRAY, italic=True)
disclaimer(s)

# ---------------- Slide 10: Demo + Contact ----------------
s = slide(); bg(s, CHARCOAL)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SH)
band.fill.solid(); band.fill.fore_color.rgb = AARP_RED; band.line.fill.background(); band.shadow.inherit = False
_, tf = box(s, Inches(1.0), Inches(0.75), Inches(11.0), Inches(0.4))
setpara(tf.paragraphs[0], "SEE IT LIVE  ·  NEXT STEP", 13, RGBColor(0xF3,0x9A,0xA6), bold=True)
_, tf = box(s, Inches(1.0), Inches(1.5), Inches(11.3), Inches(1.4))
setpara(tf.paragraphs[0], "Let's open the AARP-branded demo", 40, WHITE, bold=True)
_, tf = box(s, Inches(1.0), Inches(2.85), Inches(11.0), Inches(1.6))
bullets(tf, [
    ("Request a private demo.", "The white-label build is kept private as a partner concept — we'll provide guided access at parigrado.com/?partner=aarp."),
    ("30-minute intro call.", "Align on pilot states, success metrics, and scope; live embed within two weeks of a signed pilot."),
], size=16, gap=12)
# override bullet text colors to light for dark bg
for para in tf.paragraphs:
    for r in para.runs:
        if r.font.color and r.font.color.type is not None and r.font.color.rgb == CHARCOAL:
            r.font.color.rgb = WHITE
        elif r.font.color and r.font.color.type is not None and r.font.color.rgb == GRAY:
            r.font.color.rgb = RGBColor(0xC9,0xCC,0xD1)
# contact card
rect(s, Inches(1.0), Inches(5.15), Inches(11.3), Inches(1.5), fill=RGBColor(0x3A,0x3A,0x3A))
_, tf = box(s, Inches(1.4), Inches(5.35), Inches(10.6), Inches(1.2))
setpara(tf.paragraphs[0], "Larry Kirschner  ·  Parigrado", 22, WHITE, bold=True)
p = tf.add_paragraph()
r = p.add_run(); r.text = "LarryRkirschner@gmail.com    ·    239-304-6180"
r.font.size = Pt(18); r.font.name = FONT; r.font.color.rgb = RGBColor(0xF3,0x9A,0xA6); r.font.bold = True
p.space_before = Pt(6)
_, tf = box(s, Inches(1.0), Inches(7.0), Inches(11.3), Inches(0.4))
setpara(tf.paragraphs[0],
        "Partnership concept — not an official AARP product.  parigrado.com", 10, RGBColor(0x9A,0x9E,0xA5), italic=True)

out = "/Users/beta/Dropbox/hospital-compare/docs/AARP-Parigrado-Pitch-Deck.pptx"
prs.save(out)
print("Saved", out, "slides:", len(prs.slides._sldIdLst))
